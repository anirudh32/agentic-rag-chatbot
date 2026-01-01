from mcp import send_mcp_message
import fitz  
from pptx import Presentation
import pandas as pd
from docx import Document
import logging
from langchain.text_splitter import RecursiveCharacterTextSplitter

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=200,
    separators=["\n\n", "\n", ".", " ", ""]
)

# Parse PDF using PyMuPDF
def parse_pdf(file):
    try:
        doc = fitz.open(stream=file.read(), filetype="pdf")
        text = ""
        for page in doc:
            text += page.get_text("text") or ""
        doc.close()
        return text
    except Exception as e:
        logging.error(f"Error parsing PDF: {str(e)}")
        return None

# Parse PPTX
def parse_pptx(file):
    try:
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logging.error(f"Error parsing PPTX: {str(e)}")
        return None

# Parse CSV
def parse_csv(file):
    try:
        df = pd.read_csv(file)
        return df.to_string()
    except Exception as e:
        logging.error(f"Error parsing CSV: {str(e)}")
        return None

# Parse DOCX
def parse_docx(file):
    try:
        doc = Document(file)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    except Exception as e:
        logging.error(f"Error parsing DOCX: {str(e)}")
        return None

# Parse TXT/Markdown
def parse_txt(file):
    try:
        return file.read().decode("utf-8")
    except Exception as e:
        logging.error(f"Error parsing TXT: {str(e)}")
        return None

# IngestionAgent: Process document
def process_document(file, filename):
    extension = filename.lower().split(".")[-1]
    if extension == "pdf":
        text = parse_pdf(file)
    elif extension == "pptx":
        text = parse_pptx(file)
    elif extension == "csv":
        text = parse_csv(file)
    elif extension == "docx":
        text = parse_docx(file)
    elif extension in ["txt", "md"]:
        text = parse_txt(file)
    else:
        return False

    if text:
        # Split text into chunks using RecursiveCharacterTextSplitter
        chunks = text_splitter.split_text(text)
        send_mcp_message(
            "IngestionAgent", "RetrievalAgent", "DOCUMENT_PARSED",
            {"chunks": chunks, "document_name": filename}
        )
        return True
    return False